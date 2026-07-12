from io import BytesIO

from fastapi import HTTPException
from pptx import Presentation

from app.config import settings


def parse_pptx(file_bytes: bytes) -> list[dict]:
    prs = Presentation(BytesIO(file_bytes))
    slides = list(prs.slides)
    if len(slides) > settings.max_pptx_slides:
        raise HTTPException(
            status_code=413,
            detail=f"Presentation has {len(slides)} slides; the limit is {settings.max_pptx_slides}.",
        )

    results = []
    for i, slide in enumerate(slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    texts.append(" | ".join(cell.text for cell in row.cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                texts.append(f"Speaker notes: {notes}")

        combined = "\n".join(texts).strip()
        if combined:
            results.append({"text": combined, "metadata": {"slide": i + 1}})

    return results
